"""
Google Sheets and Drive tools for the local agent.

Sheets tools (contact-focused):
  sheets_search_contact  — fuzzy name search in your contact sheet
  sheets_add_contact     — add a new contact row
  sheets_read            — read any range from any sheet (general purpose)

Drive tools:
  drive_search           — search for files by name
  drive_read             — read content of a Google Doc or Sheet

Auth:
  Uses credentials.json (same as Gmail/Calendar) with a separate
  token file (sheets_token.json) for Sheets+Drive scopes.
  First run opens a browser for one-time permission grant.

Config (set in .env):
  CONTACT_SHEET_ID  — the Google Sheet ID from the URL
                      (docs.google.com/spreadsheets/d/<THIS_PART>/edit)

Sheet format (row 1 = headers):
  Name | Email | Phone | Company | Relationship | Notes
"""

import os
from pathlib import Path

try:
    from google.oauth2.credentials import Credentials
    from google.auth.transport.requests import Request
    from google_auth_oauthlib.flow import InstalledAppFlow
    from googleapiclient.discovery import build
    GOOGLE_LIBS_AVAILABLE = True
except ImportError:
    GOOGLE_LIBS_AVAILABLE = False

BASE_DIR         = Path(__file__).parent.parent
CREDENTIALS_PATH = BASE_DIR / "credentials.json"
TOKEN_PATH       = BASE_DIR / "sheets_token.json"
SCOPES           = [
    "https://www.googleapis.com/auth/spreadsheets",
    "https://www.googleapis.com/auth/drive",          # full Drive — needed for move/organize
]

CONTACT_SHEET_ID = os.environ.get("CONTACT_SHEET_ID", "")
CONTACT_RANGE    = "Sheet1!A:F"   # Name|Email|Phone|Company|Relationship|Notes


# ── Auth ─────────────────────────────────────────────────────────────────────

def _get_sheets_service():
    return _get_service("sheets", "v4")

def _get_drive_service():
    return _get_service("drive", "v3")

def _get_service(api: str, version: str):
    if not GOOGLE_LIBS_AVAILABLE:
        raise RuntimeError(
            "Google API libraries not installed.\n"
            "Run: pip install google-auth google-auth-oauthlib "
            "google-auth-httplib2 google-api-python-client --break-system-packages"
        )
    if not CREDENTIALS_PATH.exists():
        raise RuntimeError(f"credentials.json not found at {CREDENTIALS_PATH}")

    creds = None
    if TOKEN_PATH.exists():
        creds = Credentials.from_authorized_user_file(str(TOKEN_PATH), SCOPES)

    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            creds.refresh(Request())
        else:
            flow = InstalledAppFlow.from_client_secrets_file(
                str(CREDENTIALS_PATH), SCOPES
            )
            creds = flow.run_local_server(port=0)
        TOKEN_PATH.write_text(creds.to_json())

    return build(api, version, credentials=creds)


# ── Helpers ───────────────────────────────────────────────────────────────────

def _sheet_to_dicts(values: list[list]) -> list[dict]:
    """Convert sheet rows (first row = headers) to list of dicts."""
    if not values or len(values) < 2:
        return []
    headers = [h.strip() for h in values[0]]
    rows = []
    for i, row in enumerate(values[1:], start=2):
        # Pad short rows with empty strings
        padded = row + [""] * (len(headers) - len(row))
        rows.append({"_row": i, **dict(zip(headers, padded))})
    return rows


def _format_contact(c: dict) -> str:
    """Format a contact dict for display."""
    parts = []
    for field in ("Name", "Email", "Phone", "Company", "Relationship", "Notes"):
        val = c.get(field, "").strip()
        if val:
            parts.append(f"  {field}: {val}")
    return "\n".join(parts)


def _fuzzy_match(query: str, name: str) -> bool:
    """Return True if query matches the name (case-insensitive, partial)."""
    q = query.lower().strip()
    n = name.lower().strip()
    # Full match, starts with, or all query words appear in name
    return (q == n or n.startswith(q) or q in n or
            all(w in n for w in q.split()))


# ── Sheets contact tools ──────────────────────────────────────────────────────

def sheets_search_contact(args: dict) -> str:
    """
    Search for contacts by name in your Google Sheets contact list.
    Returns all matching rows with full details for disambiguation.
    Args:
      name       (str) — name to search for (partial names work)
      sheet_id   (str, optional) — override CONTACT_SHEET_ID from .env
    """
    query    = (args.get("name") or "").strip()
    sheet_id = (args.get("sheet_id") or CONTACT_SHEET_ID).strip()

    if not query:
        return "Error: 'name' argument is required"
    if not sheet_id:
        return (
            "Error: CONTACT_SHEET_ID not set.\n"
            "Add it to your .env file:\n"
            "  CONTACT_SHEET_ID=your_sheet_id_here\n"
            "(get it from the Google Sheets URL)"
        )

    try:
        service = _get_sheets_service()
        result  = service.spreadsheets().values().get(
            spreadsheetId=sheet_id,
            range=CONTACT_RANGE,
        ).execute()

        rows = _sheet_to_dicts(result.get("values", []))
        if not rows:
            return "Contact sheet is empty or has no data rows."

        matches = [r for r in rows if _fuzzy_match(query, r.get("Name", ""))]

        if not matches:
            return f"No contacts found matching '{query}'."

        if len(matches) == 1:
            c = matches[0]
            return (
                f"Found 1 contact matching '{query}':\n\n"
                f"{_format_contact(c)}"
            )

        # Multiple matches — return all with enough info to disambiguate
        lines = [f"Found {len(matches)} contacts matching '{query}' — which one?\n"]
        for i, c in enumerate(matches, 1):
            name     = c.get("Name", "?")
            company  = c.get("Company", "")
            rel      = c.get("Relationship", "")
            email    = c.get("Email", "")
            disambig = " | ".join(filter(None, [company, rel, email]))
            lines.append(f"{i}. {name}  ({disambig})")

        lines.append(
            "\nPlease specify which one — e.g. 'the one from River Tech' "
            "or 'Vansh Sahni' (full name)."
        )
        return "\n".join(lines)

    except Exception as e:
        return f"Error searching contacts: {e}"


def sheets_add_contact(args: dict) -> str:
    """
    Add a new contact to the Google Sheets contact list.
    Args:
      name         (str) — full name
      email        (str) — email address
      phone        (str, optional)
      company      (str, optional)
      relationship (str, optional) — e.g. Client, Friend, Business Partner
      notes        (str, optional)
      sheet_id     (str, optional) — override CONTACT_SHEET_ID from .env
    """
    name         = (args.get("name") or "").strip()
    email        = (args.get("email") or "").strip()
    phone        = (args.get("phone") or "").strip()
    company      = (args.get("company") or "").strip()
    relationship = (args.get("relationship") or "").strip()
    notes        = (args.get("notes") or "").strip()
    sheet_id     = (args.get("sheet_id") or CONTACT_SHEET_ID).strip()

    if not name:
        return "Error: 'name' is required"
    if not email:
        return "Error: 'email' is required"
    if not sheet_id:
        return "Error: CONTACT_SHEET_ID not set in .env"

    try:
        service = _get_sheets_service()
        service.spreadsheets().values().append(
            spreadsheetId=sheet_id,
            range=CONTACT_RANGE,
            valueInputOption="USER_ENTERED",
            insertDataOption="INSERT_ROWS",
            body={"values": [[name, email, phone, company, relationship, notes]]},
        ).execute()
        return (
            f"Contact added successfully!\n"
            f"  Name:         {name}\n"
            f"  Email:        {email}\n"
            f"  Phone:        {phone or '-'}\n"
            f"  Company:      {company or '-'}\n"
            f"  Relationship: {relationship or '-'}\n"
            f"  Notes:        {notes or '-'}"
        )
    except Exception as e:
        return f"Error adding contact: {e}"


def sheets_read(args: dict) -> str:
    """
    Read a range from any Google Sheet.
    Args:
      sheet_id  (str) — the Google Sheet ID
      range     (str, optional) — A1 notation range e.g. 'Sheet1!A1:E20' (default: full Sheet1)
      as_table  (bool, optional) — format as readable table (default true)
    """
    sheet_id = (args.get("sheet_id") or CONTACT_SHEET_ID).strip()
    rng      = (args.get("range") or "Sheet1").strip()
    as_table = str(args.get("as_table", "true")).lower() != "false"

    if not sheet_id:
        return "Error: 'sheet_id' is required"

    try:
        service = _get_sheets_service()
        result  = service.spreadsheets().values().get(
            spreadsheetId=sheet_id,
            range=rng,
        ).execute()

        values = result.get("values", [])
        if not values:
            return "No data found in that range."

        if not as_table or len(values) < 2:
            # Raw output
            return "\n".join("\t".join(str(c) for c in row) for row in values)

        # Table output with headers
        headers = values[0]
        rows    = values[1:]
        col_widths = [max(len(str(h)), max((len(str(r[i])) if i < len(r) else 0) for r in rows), 1)
                      for i, h in enumerate(headers)]

        def fmt_row(row):
            return " | ".join(
                str(row[i] if i < len(row) else "").ljust(col_widths[i])
                for i in range(len(headers))
            )

        sep = "-+-".join("-" * w for w in col_widths)
        lines = [fmt_row(headers), sep] + [fmt_row(r) for r in rows]
        return f"Sheet data ({len(rows)} rows):\n\n" + "\n".join(lines)

    except Exception as e:
        return f"Error reading sheet: {e}"


# ── Drive tools ───────────────────────────────────────────────────────────────

def drive_list(args: dict) -> str:
    """
    List files and folders in Google Drive.
    Args:
      folder_id (str, optional) — list contents of a specific folder ID
                                   (default: root / My Drive)
      max       (int, optional, default 20) — max files to return
      type      (str, optional) — filter: 'doc', 'sheet', 'pdf', 'folder'
    """
    folder_id = (args.get("folder_id") or "root").strip()
    max_res   = int(args.get("max", 20))
    filetype  = (args.get("type") or "").lower()

    mime_map = {
        "doc":    "application/vnd.google-apps.document",
        "sheet":  "application/vnd.google-apps.spreadsheet",
        "pdf":    "application/pdf",
        "folder": "application/vnd.google-apps.folder",
    }

    q_parts = [f"'{folder_id}' in parents", "trashed = false"]
    if filetype in mime_map:
        q_parts.append(f"mimeType = '{mime_map[filetype]}'")

    try:
        service = _get_drive_service()
        result  = service.files().list(
            q=" and ".join(q_parts),
            pageSize=max_res,
            orderBy="folder,name",
            fields="files(id, name, mimeType, modifiedTime, size, webViewLink)",
        ).execute()

        files = result.get("files", [])
        if not files:
            return "No files found in Google Drive (or Drive is empty)."

        lines = [f"Google Drive contents ({len(files)} items):\n"]
        for f in files:
            mime  = f.get("mimeType", "")
            # Human-readable type label
            if "folder" in mime:
                icon = "📁"
                type_label = "Folder"
            elif "document" in mime:
                icon = "📝"
                type_label = "Google Doc"
            elif "spreadsheet" in mime:
                icon = "📊"
                type_label = "Google Sheet"
            elif "presentation" in mime:
                icon = "📊"
                type_label = "Google Slides"
            elif "pdf" in mime:
                icon = "📄"
                type_label = "PDF"
            else:
                icon = "📎"
                type_label = mime.split("/")[-1]

            mtime = f.get("modifiedTime", "")[:10]
            lines.append(
                f"{icon} {f['name']}\n"
                f"   Type: {type_label}  |  Modified: {mtime}\n"
                f"   ID: {f['id']}"
            )
        return "\n".join(lines)

    except Exception as e:
        return f"Error listing Drive files: {e}"



    """
    Search for files in Google Drive by name.
    Args:
      query    (str) — file name or partial name to search for
      max      (int, optional, default 10) — max results
      type     (str, optional) — filter by type: 'doc', 'sheet', 'pdf', 'folder'
    """
    query    = (args.get("query") or "").strip()
    max_res  = int(args.get("max", 10))
    filetype = (args.get("type") or "").lower()

    if not query:
        return "Error: 'query' argument is required"

    mime_map = {
        "doc":    "application/vnd.google-apps.document",
        "sheet":  "application/vnd.google-apps.spreadsheet",
        "pdf":    "application/pdf",
        "folder": "application/vnd.google-apps.folder",
    }

    q_parts = [f"name contains '{query}'", "trashed = false"]
    if filetype in mime_map:
        q_parts.append(f"mimeType = '{mime_map[filetype]}'")

    try:
        service = _get_drive_service()
        result  = service.files().list(
            q=" and ".join(q_parts),
            pageSize=max_res,
            fields="files(id, name, mimeType, modifiedTime, size, webViewLink)",
        ).execute()

        files = result.get("files", [])
        if not files:
            return f"No files found matching '{query}'."

        lines = [f"Found {len(files)} file(s) matching '{query}':\n"]
        for f in files:
            mime  = f.get("mimeType", "").split(".")[-1].replace("google-apps.", "")
            mtime = f.get("modifiedTime", "")[:10]
            lines.append(
                f"• {f['name']}\n"
                f"  ID: {f['id']}\n"
                f"  Type: {mime}  |  Modified: {mtime}\n"
                f"  Link: {f.get('webViewLink', '-')}"
            )
        return "\n".join(lines)

    except Exception as e:
        return f"Error searching Drive: {e}"


def drive_search(args: dict) -> str:
    """
    Search for files in Google Drive by name.
    Args:
      query    (str) — file name or partial name to search for
      max      (int, optional, default 10) — max results
      type     (str, optional) — filter by type: 'doc', 'sheet', 'slides', 'pdf', 'folder'
    """
    query    = (args.get("query") or "").strip()
    max_res  = int(args.get("max", 10))
    filetype = (args.get("type") or "").lower()

    if not query:
        return "Error: 'query' argument is required"

    mime_map = {
        "doc":    "application/vnd.google-apps.document",
        "sheet":  "application/vnd.google-apps.spreadsheet",
        "slides": "application/vnd.google-apps.presentation",
        "pdf":    "application/pdf",
        "folder": "application/vnd.google-apps.folder",
    }

    q_parts = [f"name contains '{query}'", "trashed = false"]
    if filetype in mime_map:
        q_parts.append(f"mimeType = '{mime_map[filetype]}'")

    try:
        service = _get_drive_service()
        result  = service.files().list(
            q=" and ".join(q_parts),
            pageSize=max_res,
            fields="files(id, name, mimeType, modifiedTime, size, webViewLink)",
        ).execute()

        files = result.get("files", [])
        if not files:
            return f"No files found matching '{query}'."

        lines = [f"Found {len(files)} file(s) matching '{query}':\n"]
        for f in files:
            mime  = f.get("mimeType", "").split(".")[-1].replace("google-apps.", "")
            mtime = f.get("modifiedTime", "")[:10]
            lines.append(
                f"• {f['name']}\n"
                f"  ID: {f['id']}\n"
                f"  Type: {mime}  |  Modified: {mtime}\n"
                f"  Link: {f.get('webViewLink', '-')}"
            )
        return "\n".join(lines)

    except Exception as e:
        return f"Error searching Drive: {e}"


def drive_read(args: dict) -> str:
    """
    Read the text content of a file from Google Drive. Supports Google Docs,
    Sheets, Slides, plain text/markdown files, and PDFs (text extraction).
    Args:
      file_id   (str) — the Drive file ID (from drive_list/drive_search)
      max_chars (int, optional, default 6000) — max characters to return
    """
    file_id   = (args.get("file_id") or "").strip()
    max_chars = int(args.get("max_chars", 6000))

    if not file_id:
        return "Error: 'file_id' is required — get it from drive_list or drive_search"

    try:
        service = _get_drive_service()

        meta = service.files().get(
            fileId=file_id, fields="name,mimeType"
        ).execute()
        mime = meta.get("mimeType", "")
        name = meta.get("name", file_id)
        text = None

        # ── Google-native formats: exportable as plain text ──
        if "google-apps.document" in mime:
            resp = service.files().export(fileId=file_id, mimeType="text/plain").execute()
            text = resp.decode("utf-8") if isinstance(resp, bytes) else str(resp)

        elif "google-apps.spreadsheet" in mime:
            resp = service.files().export(fileId=file_id, mimeType="text/csv").execute()
            text = resp.decode("utf-8") if isinstance(resp, bytes) else str(resp)

        elif "google-apps.presentation" in mime:
            resp = service.files().export(fileId=file_id, mimeType="text/plain").execute()
            text = resp.decode("utf-8") if isinstance(resp, bytes) else str(resp)

        elif "google-apps.form" in mime:
            return (
                f"'{name}' is a Google Form. Form content (questions/responses) "
                "isn't readable via Drive export — open it directly in Google Forms "
                "to view or edit it."
            )

        # ── Plain text-like files: download directly ──
        elif mime in ("text/plain", "text/markdown", "text/csv", "application/json"):
            raw = service.files().get_media(fileId=file_id).execute()
            text = raw.decode("utf-8", errors="replace") if isinstance(raw, bytes) else str(raw)

        # ── PDF: download and extract text if pypdf is available ──
        elif mime == "application/pdf":
            raw = service.files().get_media(fileId=file_id).execute()
            try:
                import io
                from pypdf import PdfReader
                reader = PdfReader(io.BytesIO(raw))
                pages = [p.extract_text() or "" for p in reader.pages]
                text = "\n\n".join(pages).strip()
                if not text:
                    return f"'{name}' is a PDF but no extractable text was found (may be scanned/image-based)."
            except ImportError:
                return (
                    f"'{name}' is a PDF. To read PDF content, install pypdf:\n"
                    "  pip install pypdf --break-system-packages"
                )

        # ── Microsoft Office formats — parsed natively, no conversion needed ──
        elif "wordprocessingml" in mime or mime == "application/msword":
            raw = service.files().get_media(fileId=file_id).execute()
            text = _extract_docx_text(raw, name)

        elif "presentationml" in mime or mime == "application/vnd.ms-powerpoint":
            raw = service.files().get_media(fileId=file_id).execute()
            text = _extract_pptx_text(raw, name)

        elif "spreadsheetml" in mime or mime == "application/vnd.ms-excel":
            raw = service.files().get_media(fileId=file_id).execute()
            text = _extract_xlsx_text(raw, name)

        else:
            return f"'{name}' is a binary file ({mime}) — can't extract readable text from this format."

        if text is None:
            return f"Could not extract text from '{name}'."
        if isinstance(text, str) and text.startswith("__ERROR__:"):
            return text.replace("__ERROR__:", "").strip()

        text = text.strip()
        if len(text) > max_chars:
            text = text[:max_chars] + f"\n\n[... truncated at {max_chars} chars]"

        return f"Content of '{name}':\n\n{text}"

    except Exception as e:
        return f"Error reading Drive file: {e}"


def _extract_docx_text(raw_bytes: bytes, name: str) -> str:
    """Extract all text (paragraphs + tables) from a .docx file."""
    try:
        import io
        from docx import Document
    except ImportError:
        return (
            f"__ERROR__:'{name}' is a Word document. To read it, install python-docx:\n"
            "  pip install python-docx --break-system-packages"
        )

    try:
        doc = Document(io.BytesIO(raw_bytes))
        parts = []

        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)

        # Tables often hold structured data like invoice line items
        for t_idx, table in enumerate(doc.tables, 1):
            parts.append(f"\n[Table {t_idx}]")
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells]
                parts.append(" | ".join(cells))

        return "\n".join(parts)
    except Exception as e:
        return f"__ERROR__:Error parsing .docx file '{name}': {e}"


def _extract_pptx_text(raw_bytes: bytes, name: str) -> str:
    """Extract all text from a .pptx file, organized by slide."""
    try:
        import io
        from pptx import Presentation
    except ImportError:
        return (
            f"__ERROR__:'{name}' is a PowerPoint file. To read it, install python-pptx:\n"
            "  pip install python-pptx --break-system-packages"
        )

    try:
        prs = Presentation(io.BytesIO(raw_bytes))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = "".join(run.text for run in para.runs)
                        if line.strip():
                            slide_text.append(line)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells]
                        slide_text.append(" | ".join(cells))
            if slide_text:
                parts.append(f"[Slide {i}]\n" + "\n".join(slide_text))
        return "\n\n".join(parts)
    except Exception as e:
        return f"__ERROR__:Error parsing .pptx file '{name}': {e}"


def _extract_xlsx_text(raw_bytes: bytes, name: str) -> str:
    """Extract all cell data from a .xlsx file as readable tab-separated text."""
    try:
        import io
        from openpyxl import load_workbook
    except ImportError:
        return (
            f"__ERROR__:'{name}' is an Excel file. To read it, install openpyxl:\n"
            "  pip install openpyxl --break-system-packages"
        )

    try:
        wb = load_workbook(io.BytesIO(raw_bytes), data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"[Sheet: {sheet_name}]")
            for row in ws.iter_rows(values_only=True):
                if any(c is not None for c in row):
                    parts.append(" | ".join(str(c) if c is not None else "" for c in row))
        return "\n".join(parts)
    except Exception as e:
        return f"__ERROR__:Error parsing .xlsx file '{name}': {e}"

# ── Sheets creation and editing ───────────────────────────────────────────────

def sheets_create(args: dict) -> str:
    """
    Create a new Google Spreadsheet.
    Args:
      name     (str) — spreadsheet title
      headers  (list or str, optional) — column headers for row 1
                e.g. ["Name", "Revenue", "Date"] or "Name,Revenue,Date"
      data     (list, optional) — initial rows of data (list of lists)
    Returns the new spreadsheet ID and URL.
    """
    import json

    name        = (args.get("name") or "").strip()
    headers     = args.get("headers")
    data        = args.get("data")
    folder_name = (args.get("folder_name") or "").strip()
    folder_id   = (args.get("folder_id") or "").strip()

    if not name:
        return "Error: 'name' is required"

    # Parse headers — accept list or comma-separated string
    header_row = []
    if isinstance(headers, str) and headers.strip():
        header_row = [h.strip() for h in headers.split(",")]
    elif isinstance(headers, list):
        header_row = [str(h) for h in headers]

    # Parse data — accept list of lists or JSON string
    data_rows = []
    if isinstance(data, str) and data.strip():
        try:
            data_rows = json.loads(data)
        except json.JSONDecodeError:
            return "Error: 'data' must be a JSON array of arrays e.g. [[\"a\",1],[\"b\",2]]"
    elif isinstance(data, list):
        data_rows = data

    try:
        service = _get_sheets_service()

        # Resolve folder by name if given
        if folder_name and not folder_id:
            drive_svc = _get_drive_service()
            result = drive_svc.files().list(
                q=f"name = '{folder_name}' and mimeType = 'application/vnd.google-apps.folder' and trashed = false",
                fields="files(id, name)", pageSize=3,
            ).execute()
            folders = result.get("files", [])
            if not folders:
                return f"Folder '{folder_name}' not found. Create it in Drive first or use drive_move after creation."
            if len(folders) > 1:
                opts = "\n".join(f"• {f['name']} (ID: {f['id']})" for f in folders)
                return f"Multiple folders match '{folder_name}':\n{opts}\nPlease specify folder_id directly."
            folder_id = folders[0]["id"]

        # Build spreadsheet body with optional parent folder
        create_body: dict = {"properties": {"title": name}}
        if folder_id:
            # We create via Drive API to set parent folder
            drive_svc = _get_drive_service()
            drive_meta = {
                "name": name,
                "mimeType": "application/vnd.google-apps.spreadsheet",
                "parents": [folder_id],
            }
            created_file = drive_svc.files().create(
                body=drive_meta, fields="id, webViewLink"
            ).execute()
            sheet_id  = created_file["id"]
            sheet_url = created_file.get("webViewLink", f"https://docs.google.com/spreadsheets/d/{sheet_id}")
        else:
            # Create via Sheets API (lands in root My Drive)
            spreadsheet = service.spreadsheets().create(
                body=create_body,
                fields="spreadsheetId,spreadsheetUrl",
            ).execute()
            sheet_id  = spreadsheet["spreadsheetId"]
            sheet_url = spreadsheet["spreadsheetUrl"]

        # Write headers and initial data if provided
        all_rows = []
        if header_row:
            all_rows.append(header_row)
        if data_rows:
            all_rows.extend(data_rows)

        if all_rows:
            service.spreadsheets().values().update(
                spreadsheetId=sheet_id,
                range="Sheet1!A1",
                valueInputOption="USER_ENTERED",
                body={"values": all_rows},
            ).execute()

        result = (
            f"Spreadsheet created: '{name}'\n"
            f"ID:  {sheet_id}\n"
            f"URL: {sheet_url}\n"
        )
        if header_row:
            result += f"Headers: {', '.join(header_row)}\n"
        if data_rows:
            result += f"Initial rows written: {len(data_rows)}"

        return result

    except Exception as e:
        return f"Error creating spreadsheet: {e}"


def sheets_write(args: dict) -> str:
    """
    Write (overwrite) data to a specific range in any Google Sheet.
    Use this to set cell values, update a table, or write headers.
    Args:
      sheet_id  (str) — spreadsheet ID (from sheets_create or Drive URL)
      range     (str) — A1 notation e.g. 'Sheet1!A1' or 'Sheet1!B2:D5'
      data      (list or str) — 2D array of values:
                  [["Name","Age"],["Daksh",18]] or JSON string of same
      clear_first (bool, optional) — clear the range before writing (default false)
    """
    import json

    sheet_id    = (args.get("sheet_id") or "").strip()
    rng         = (args.get("range") or "Sheet1!A1").strip()
    data        = args.get("data")
    clear_first = str(args.get("clear_first", "false")).lower() == "true"

    if not sheet_id:
        return "Error: 'sheet_id' is required"
    if data is None:
        return "Error: 'data' is required"

    # Parse data
    if isinstance(data, str):
        try:
            rows = json.loads(data)
        except json.JSONDecodeError:
            # Try treating as a single row of comma-separated values
            rows = [[v.strip() for v in data.split(",")]]
    elif isinstance(data, list):
        rows = data
    else:
        return f"Error: 'data' must be a list or JSON string, got {type(data).__name__}"

    # Ensure rows is a list of lists
    if rows and not isinstance(rows[0], list):
        rows = [rows]  # wrap single row

    try:
        service = _get_sheets_service()

        if clear_first:
            service.spreadsheets().values().clear(
                spreadsheetId=sheet_id, range=rng
            ).execute()

        result = service.spreadsheets().values().update(
            spreadsheetId=sheet_id,
            range=rng,
            valueInputOption="USER_ENTERED",
            body={"values": rows},
        ).execute()

        updated_cells = result.get("updatedCells", "?")
        updated_range = result.get("updatedRange", rng)
        return (
            f"Written successfully to '{updated_range}'\n"
            f"Rows written: {len(rows)}  |  Cells updated: {updated_cells}\n"
            f"Sheet ID: {sheet_id}"
        )

    except Exception as e:
        return f"Error writing to sheet: {e}"


def sheets_append(args: dict) -> str:
    """
    Append new rows to the end of existing data in a Google Sheet.
    Finds the last row with data and adds below it — safe to use repeatedly.
    Args:
      sheet_id  (str) — spreadsheet ID
      data      (list or str) — rows to append:
                  [["Daksh",18],["Vansh",20]] or JSON string of same
      sheet_name (str, optional) — sheet/tab name (default: Sheet1)
    """
    import json

    sheet_id   = (args.get("sheet_id") or "").strip()
    data       = args.get("data")
    sheet_name = (args.get("sheet_name") or "Sheet1").strip()

    if not sheet_id:
        return "Error: 'sheet_id' is required"
    if data is None:
        return "Error: 'data' is required"

    # Parse data
    if isinstance(data, str):
        try:
            rows = json.loads(data)
        except json.JSONDecodeError:
            rows = [[v.strip() for v in data.split(",")]]
    elif isinstance(data, list):
        rows = data
    else:
        return f"Error: 'data' must be a list or JSON string"

    if rows and not isinstance(rows[0], list):
        rows = [rows]

    try:
        service = _get_sheets_service()
        result  = service.spreadsheets().values().append(
            spreadsheetId=sheet_id,
            range=f"{sheet_name}!A1",
            valueInputOption="USER_ENTERED",
            insertDataOption="INSERT_ROWS",
            body={"values": rows},
        ).execute()

        updated_range = result.get("updates", {}).get("updatedRange", "?")
        return (
            f"Appended {len(rows)} row(s) to '{sheet_name}'\n"
            f"Written to: {updated_range}\n"
            f"Sheet ID: {sheet_id}"
        )

    except Exception as e:
        return f"Error appending to sheet: {e}"

# ── Drive file organization ───────────────────────────────────────────────────

def drive_move(args: dict) -> str:
    """
    Move a file or folder to a different folder in Google Drive.
    Args:
      file_id      (str) — ID of the file to move (from drive_list/drive_search)
      folder_name  (str) — name of the destination folder to search for
      folder_id    (str) — destination folder ID (use instead of folder_name
                           if you already have the ID from drive_list)
    If folder_name matches multiple folders, returns options to disambiguate.
    """
    file_id     = (args.get("file_id") or "").strip()
    folder_name = (args.get("folder_name") or "").strip()
    folder_id   = (args.get("folder_id") or "").strip()

    if not file_id:
        return "Error: 'file_id' is required — get it from drive_list or drive_search"
    if not folder_name and not folder_id:
        return "Error: either 'folder_name' or 'folder_id' is required"

    try:
        service = _get_drive_service()

        # If folder_name given, search for matching folders
        if folder_name and not folder_id:
            result = service.files().list(
                q=f"name = '{folder_name}' and mimeType = 'application/vnd.google-apps.folder' and trashed = false",
                fields="files(id, name)",
                pageSize=5,
            ).execute()
            folders = result.get("files", [])

            if not folders:
                # Try partial match
                result = service.files().list(
                    q=f"name contains '{folder_name}' and mimeType = 'application/vnd.google-apps.folder' and trashed = false",
                    fields="files(id, name)",
                    pageSize=5,
                ).execute()
                folders = result.get("files", [])

            if not folders:
                return f"No folder found matching '{folder_name}'. Use drive_list with type='folder' to see your folders."

            if len(folders) > 1:
                lines = [f"Multiple folders match '{folder_name}' — which one?\n"]
                for f in folders:
                    lines.append(f"• {f['name']}  (ID: {f['id']})")
                lines.append("\nReply with the folder ID or a more specific name.")
                return "\n".join(lines)

            folder_id = folders[0]["id"]
            resolved_name = folders[0]["name"]
        else:
            resolved_name = folder_id  # use ID as display if no name given

        # Get the file's current metadata (name + current parent)
        file_meta = service.files().get(
            fileId=file_id,
            fields="name, parents",
        ).execute()
        file_name      = file_meta.get("name", file_id)
        current_parents = ",".join(file_meta.get("parents", []))

        # Move: add new parent, remove old parent(s)
        service.files().update(
            fileId=file_id,
            addParents=folder_id,
            removeParents=current_parents,
            fields="id, parents",
        ).execute()

        return (
            f"Moved '{file_name}' to '{resolved_name}' successfully.\n"
            f"File ID:   {file_id}\n"
            f"Folder ID: {folder_id}"
        )

    except Exception as e:
        return f"Error moving file: {e}"

# ── Document understanding (analyze, don't just dump text) ───────────────────

_llm_client = None

def set_llm_client(client) -> None:
    """Wire in the LLM client so drive_read_and_explain can analyze content."""
    global _llm_client
    _llm_client = client


def drive_read_and_explain(args: dict) -> str:
    """
    Read a Drive file AND have the LLM analyze/explain its contents —
    e.g. count invoices, extract names and amounts, summarize key points.
    Use this instead of drive_read when Daksh asks "what's inside" or
    "tell me about" a document, rather than wanting the raw text.
    Args:
      file_id   (str) — the Drive file ID (from drive_list/drive_search)
      question  (str, optional) — specific question to answer about the content
                                   e.g. "how many invoices and what are the amounts?"
    """
    file_id  = (args.get("file_id") or "").strip()
    question = (args.get("question") or "").strip()

    if not file_id:
        return "Error: 'file_id' is required — get it from drive_list or drive_search"
    if not _llm_client:
        return "Error: LLM client not initialized for document analysis"

    # Reuse drive_read to get the raw extracted text
    raw_result = drive_read({"file_id": file_id, "max_chars": 12000})
    if raw_result.startswith("Error") or raw_result.startswith("'") and "binary file" in raw_result:
        return raw_result  # pass through extraction errors as-is

    analysis_prompt = (
        "You are analyzing a document's extracted content. Answer clearly and "
        "concisely. If the document contains structured data (invoices, line "
        "items, tables, lists of people/amounts), extract and organize that "
        "data explicitly — counts, names, amounts, dates. Don't just describe "
        "the document in vague terms.\n\n"
        f"{raw_result}\n\n"
        f"Task: {question if question else 'Summarize what this document contains, calling out any structured data like names, amounts, dates, or counts.'}"
    )

    try:
        messages = [{"role": "user", "content": analysis_prompt}]
        analysis = _llm_client.chat(messages)
        return analysis.strip() if analysis else "Analysis returned empty response."
    except Exception as e:
        return f"Document was read successfully but analysis failed: {e}\n\nRaw content:\n{raw_result[:1000]}"
